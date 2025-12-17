"""
File format tools for working with common file types.

Provides tools for reading, writing, and manipulating:
- Excel files (.xlsx)
- CSV files
- JSON files
- PDF files (text extraction)
- PowerPoint files (.pptx)
- Images (metadata, resize, convert)

Install dependencies: pip install tool-master[files]
"""

import csv
import io
import json
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

from tool_master.schemas.tool import ParameterType, Tool, ToolParameter

# Maximum file size (50MB default)
MAX_FILE_SIZE = 50 * 1024 * 1024


def _validate_file_path(file_path: str, must_exist: bool = True) -> Path:
    """Validate and normalize a file path."""
    path = Path(file_path).resolve()

    if must_exist and not path.exists():
        raise ValueError(f"File not found: {file_path}")

    if must_exist and not path.is_file():
        raise ValueError(f"Not a file: {file_path}")

    if must_exist and path.stat().st_size > MAX_FILE_SIZE:
        size_mb = path.stat().st_size / (1024 * 1024)
        raise ValueError(f"File too large: {size_mb:.1f}MB (max {MAX_FILE_SIZE // (1024*1024)}MB)")

    return path


def _get_file_info(path: Path) -> dict:
    """Get basic file information."""
    stat = path.stat()
    return {
        "file_path": str(path),
        "file_name": path.name,
        "file_size_bytes": stat.st_size,
        "modified_time": datetime.fromtimestamp(stat.st_mtime).isoformat(),
    }


# =============================================================================
# EXCEL TOOLS
# =============================================================================

def _read_excel(
    file_path: str,
    sheet_name: Optional[str] = None,
    max_rows: int = 1000,
) -> dict:
    """Read an Excel file and return its contents."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "openpyxl package required for Excel tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)

    if not path.suffix.lower() in ('.xlsx', '.xlsm'):
        raise ValueError(f"Not an Excel file: {path.suffix}")

    wb = load_workbook(path, read_only=True, data_only=True)

    # Get sheet to read
    if sheet_name:
        if sheet_name not in wb.sheetnames:
            raise ValueError(f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}")
        ws = wb[sheet_name]
    else:
        ws = wb.active
        sheet_name = ws.title

    # Read data
    data = []
    headers = []
    row_count = 0

    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i == 0:
            headers = [str(cell) if cell is not None else "" for cell in row]
        else:
            if i > max_rows:
                break
            # Convert cells to serializable values
            row_data = []
            for cell in row:
                if cell is None:
                    row_data.append(None)
                elif isinstance(cell, datetime):
                    row_data.append(cell.isoformat())
                else:
                    row_data.append(cell)
            data.append(row_data)
        row_count = i

    wb.close()

    result = _get_file_info(path)
    result.update({
        "sheet_name": sheet_name,
        "total_sheets": len(wb.sheetnames),
        "sheet_names": wb.sheetnames,
        "row_count": row_count,
        "column_count": len(headers),
        "headers": headers,
        "data": data,
        "truncated": row_count > max_rows,
    })

    return result


def _write_excel(
    file_path: str,
    data: list,
    sheet_name: str = "Sheet1",
    headers: Optional[list] = None,
) -> dict:
    """Write data to an Excel file."""
    try:
        from openpyxl import Workbook
    except ImportError:
        raise ImportError(
            "openpyxl package required for Excel tools. "
            "Install with: pip install tool-master[files]"
        )

    path = Path(file_path).resolve()

    # Ensure .xlsx extension
    if not path.suffix.lower() == '.xlsx':
        path = path.with_suffix('.xlsx')

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name

    # Write headers if provided
    start_row = 1
    if headers:
        for col, header in enumerate(headers, 1):
            ws.cell(row=1, column=col, value=header)
        start_row = 2

    # Write data
    for row_idx, row_data in enumerate(data, start_row):
        for col_idx, value in enumerate(row_data, 1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    wb.save(path)
    wb.close()

    return {
        "file_path": str(path),
        "sheet_name": sheet_name,
        "rows_written": len(data),
        "columns": len(headers) if headers else (len(data[0]) if data else 0),
        "has_headers": headers is not None,
    }


def _list_excel_sheets(file_path: str) -> dict:
    """List all sheets in an Excel workbook."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "openpyxl package required for Excel tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)
    wb = load_workbook(path, read_only=True)

    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        sheets.append({
            "name": name,
            "is_active": name == wb.active.title,
        })

    wb.close()

    result = _get_file_info(path)
    result.update({
        "sheet_count": len(sheets),
        "sheets": sheets,
        "active_sheet": wb.active.title,
    })

    return result


def _read_excel_sheet_info(file_path: str, sheet_name: Optional[str] = None) -> dict:
    """Get detailed information about an Excel sheet."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "openpyxl package required for Excel tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)
    wb = load_workbook(path, read_only=True)

    if sheet_name:
        if sheet_name not in wb.sheetnames:
            raise ValueError(f"Sheet '{sheet_name}' not found")
        ws = wb[sheet_name]
    else:
        ws = wb.active
        sheet_name = ws.title

    # Get dimensions
    min_row = ws.min_row or 1
    max_row = ws.max_row or 1
    min_col = ws.min_column or 1
    max_col = ws.max_column or 1

    # Get first row as potential headers
    headers = []
    for row in ws.iter_rows(min_row=1, max_row=1, values_only=True):
        headers = [str(cell) if cell is not None else "" for cell in row]
        break

    wb.close()

    result = _get_file_info(path)
    result.update({
        "sheet_name": sheet_name,
        "dimensions": {
            "min_row": min_row,
            "max_row": max_row,
            "min_column": min_col,
            "max_column": max_col,
            "total_rows": max_row - min_row + 1,
            "total_columns": max_col - min_col + 1,
        },
        "headers": headers,
        "estimated_cells": (max_row - min_row + 1) * (max_col - min_col + 1),
    })

    return result


# =============================================================================
# CSV TOOLS
# =============================================================================

def _read_csv(
    file_path: str,
    delimiter: Optional[str] = None,
    max_rows: int = 1000,
    encoding: str = "utf-8",
) -> dict:
    """Read a CSV file and return its contents."""
    path = _validate_file_path(file_path)

    with open(path, 'r', encoding=encoding, newline='') as f:
        # Auto-detect delimiter if not provided
        if delimiter is None:
            sample = f.read(8192)
            f.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample)
                delimiter = dialect.delimiter
            except csv.Error:
                delimiter = ','

        reader = csv.reader(f, delimiter=delimiter)

        headers = []
        data = []
        row_count = 0

        for i, row in enumerate(reader):
            if i == 0:
                headers = row
            else:
                if i > max_rows:
                    break
                data.append(row)
            row_count = i

    result = _get_file_info(path)
    result.update({
        "delimiter": delimiter,
        "encoding": encoding,
        "row_count": row_count,
        "column_count": len(headers),
        "headers": headers,
        "data": data,
        "truncated": row_count > max_rows,
    })

    return result


def _write_csv(
    file_path: str,
    data: list,
    headers: Optional[list] = None,
    delimiter: str = ",",
    encoding: str = "utf-8",
) -> dict:
    """Write data to a CSV file."""
    path = Path(file_path).resolve()

    # Ensure .csv extension
    if not path.suffix.lower() == '.csv':
        path = path.with_suffix('.csv')

    with open(path, 'w', encoding=encoding, newline='') as f:
        writer = csv.writer(f, delimiter=delimiter)

        if headers:
            writer.writerow(headers)

        writer.writerows(data)

    return {
        "file_path": str(path),
        "rows_written": len(data),
        "columns": len(headers) if headers else (len(data[0]) if data else 0),
        "has_headers": headers is not None,
        "delimiter": delimiter,
        "encoding": encoding,
    }


def _csv_to_excel(
    csv_path: str,
    excel_path: Optional[str] = None,
    sheet_name: str = "Sheet1",
    delimiter: Optional[str] = None,
) -> dict:
    """Convert a CSV file to Excel format."""
    try:
        from openpyxl import Workbook
    except ImportError:
        raise ImportError(
            "openpyxl package required for Excel tools. "
            "Install with: pip install tool-master[files]"
        )

    # Read CSV
    csv_data = _read_csv(csv_path, delimiter=delimiter, max_rows=100000)

    # Determine Excel path
    if excel_path is None:
        excel_path = str(Path(csv_path).with_suffix('.xlsx'))

    # Write to Excel
    result = _write_excel(
        excel_path,
        csv_data["data"],
        sheet_name=sheet_name,
        headers=csv_data["headers"],
    )

    result["source_csv"] = csv_path
    result["rows_converted"] = csv_data["row_count"]

    return result


# =============================================================================
# JSON TOOLS
# =============================================================================

def _read_json(file_path: str, encoding: str = "utf-8") -> dict:
    """Read a JSON file and return its contents."""
    path = _validate_file_path(file_path)

    with open(path, 'r', encoding=encoding) as f:
        data = json.load(f)

    # Analyze structure
    def get_type_info(obj: Any) -> dict:
        if isinstance(obj, dict):
            return {"type": "object", "keys": list(obj.keys())[:20], "key_count": len(obj)}
        elif isinstance(obj, list):
            return {"type": "array", "length": len(obj), "sample_type": type(obj[0]).__name__ if obj else None}
        else:
            return {"type": type(obj).__name__}

    result = _get_file_info(path)
    result.update({
        "encoding": encoding,
        "structure": get_type_info(data),
        "data": data,
    })

    return result


def _write_json(
    file_path: str,
    data: Any,
    pretty: bool = True,
    encoding: str = "utf-8",
) -> dict:
    """Write data to a JSON file."""
    path = Path(file_path).resolve()

    # Ensure .json extension
    if not path.suffix.lower() == '.json':
        path = path.with_suffix('.json')

    with open(path, 'w', encoding=encoding) as f:
        if pretty:
            json.dump(data, f, indent=2, default=str)
        else:
            json.dump(data, f, default=str)

    return {
        "file_path": str(path),
        "file_size_bytes": path.stat().st_size,
        "pretty_printed": pretty,
        "encoding": encoding,
    }


def _validate_json(file_path: str, encoding: str = "utf-8") -> dict:
    """Validate a JSON file and return its structure."""
    path = _validate_file_path(file_path)

    try:
        with open(path, 'r', encoding=encoding) as f:
            content = f.read()
            data = json.loads(content)

        # Analyze structure recursively
        def analyze_structure(obj: Any, depth: int = 0, max_depth: int = 5) -> dict:
            if depth > max_depth:
                return {"type": "...", "truncated": True}

            if isinstance(obj, dict):
                return {
                    "type": "object",
                    "key_count": len(obj),
                    "keys": {k: analyze_structure(v, depth + 1) for k, v in list(obj.items())[:10]},
                }
            elif isinstance(obj, list):
                sample = analyze_structure(obj[0], depth + 1) if obj else None
                return {
                    "type": "array",
                    "length": len(obj),
                    "item_type": sample,
                }
            else:
                return {"type": type(obj).__name__}

        result = _get_file_info(path)
        result.update({
            "valid": True,
            "encoding": encoding,
            "structure": analyze_structure(data),
        })

        return result

    except json.JSONDecodeError as e:
        return {
            "file_path": str(path),
            "valid": False,
            "error": str(e),
            "error_line": e.lineno,
            "error_column": e.colno,
        }


# =============================================================================
# PDF TOOLS
# =============================================================================

def _read_pdf_text(
    file_path: str,
    max_pages: Optional[int] = None,
    page_numbers: Optional[list] = None,
) -> dict:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError(
            "pypdf package required for PDF tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)

    reader = PdfReader(path)
    total_pages = len(reader.pages)

    # Determine which pages to read
    if page_numbers:
        pages_to_read = [p - 1 for p in page_numbers if 0 < p <= total_pages]
    elif max_pages:
        pages_to_read = list(range(min(max_pages, total_pages)))
    else:
        pages_to_read = list(range(total_pages))

    pages = []
    for page_num in pages_to_read:
        page = reader.pages[page_num]
        text = page.extract_text() or ""
        pages.append({
            "page_number": page_num + 1,
            "text": text,
            "character_count": len(text),
        })

    result = _get_file_info(path)
    result.update({
        "total_pages": total_pages,
        "pages_extracted": len(pages),
        "pages": pages,
    })

    return result


def _read_pdf_metadata(file_path: str) -> dict:
    """Get metadata from a PDF file."""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError(
            "pypdf package required for PDF tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)

    reader = PdfReader(path)
    metadata = reader.metadata or {}

    result = _get_file_info(path)
    result.update({
        "total_pages": len(reader.pages),
        "metadata": {
            "title": metadata.get("/Title"),
            "author": metadata.get("/Author"),
            "subject": metadata.get("/Subject"),
            "creator": metadata.get("/Creator"),
            "producer": metadata.get("/Producer"),
            "creation_date": str(metadata.get("/CreationDate")) if metadata.get("/CreationDate") else None,
            "modification_date": str(metadata.get("/ModDate")) if metadata.get("/ModDate") else None,
        },
        "is_encrypted": reader.is_encrypted,
    })

    return result


def _count_pdf_pages(file_path: str) -> dict:
    """Get the page count of a PDF file."""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError(
            "pypdf package required for PDF tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)

    reader = PdfReader(path)

    result = _get_file_info(path)
    result["page_count"] = len(reader.pages)

    return result


# =============================================================================
# POWERPOINT TOOLS
# =============================================================================

def _read_pptx_text(file_path: str) -> dict:
    """Extract text from all slides in a PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx package required for PowerPoint tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)

    prs = Presentation(path)

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)

        slides.append({
            "slide_number": i,
            "text": "\n".join(texts),
            "shape_count": len(slide.shapes),
        })

    result = _get_file_info(path)
    result.update({
        "slide_count": len(slides),
        "slides": slides,
    })

    return result


def _read_pptx_structure(file_path: str) -> dict:
    """Get the structure of a PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx package required for PowerPoint tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)

    prs = Presentation(path)

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        # Try to get title
        title = None
        if slide.shapes.title:
            title = slide.shapes.title.text

        # Get notes
        notes = None
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        slides.append({
            "slide_number": i,
            "title": title,
            "has_notes": notes is not None,
            "notes_preview": notes[:200] if notes else None,
            "shape_count": len(slide.shapes),
        })

    result = _get_file_info(path)
    result.update({
        "slide_count": len(slides),
        "slides": slides,
    })

    return result


# =============================================================================
# IMAGE TOOLS
# =============================================================================

def _read_image_metadata(file_path: str) -> dict:
    """Get metadata from an image file."""
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS
    except ImportError:
        raise ImportError(
            "pillow package required for image tools. "
            "Install with: pip install tool-master[files]"
        )

    path = _validate_file_path(file_path)

    with Image.open(path) as img:
        # Basic info
        result = _get_file_info(path)
        result.update({
            "format": img.format,
            "mode": img.mode,
            "width": img.width,
            "height": img.height,
            "dimensions": f"{img.width}x{img.height}",
        })

        # EXIF data if available
        exif = {}
        if hasattr(img, '_getexif') and img._getexif():
            for tag_id, value in img._getexif().items():
                tag = TAGS.get(tag_id, tag_id)
                if isinstance(value, bytes):
                    continue  # Skip binary data
                exif[tag] = str(value) if not isinstance(value, (int, float, str)) else value

        if exif:
            result["exif"] = {k: v for k, v in list(exif.items())[:20]}  # Limit EXIF data

    return result


def _resize_image(
    file_path: str,
    width: Optional[int] = None,
    height: Optional[int] = None,
    output_path: Optional[str] = None,
    maintain_aspect: bool = True,
) -> dict:
    """Resize an image to specified dimensions."""
    try:
        from PIL import Image
    except ImportError:
        raise ImportError(
            "pillow package required for image tools. "
            "Install with: pip install tool-master[files]"
        )

    if width is None and height is None:
        raise ValueError("Must specify at least width or height")

    path = _validate_file_path(file_path)

    with Image.open(path) as img:
        original_width, original_height = img.size

        if maintain_aspect:
            if width and height:
                # Fit within box while maintaining aspect ratio
                img.thumbnail((width, height), Image.Resampling.LANCZOS)
                new_size = img.size
            elif width:
                ratio = width / original_width
                new_size = (width, int(original_height * ratio))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
            else:
                ratio = height / original_height
                new_size = (int(original_width * ratio), height)
                img = img.resize(new_size, Image.Resampling.LANCZOS)
        else:
            new_size = (width or original_width, height or original_height)
            img = img.resize(new_size, Image.Resampling.LANCZOS)

        # Determine output path
        if output_path is None:
            stem = path.stem
            output_path = str(path.with_name(f"{stem}_resized{path.suffix}"))

        output = Path(output_path).resolve()
        img.save(output)

    return {
        "input_path": str(path),
        "output_path": str(output),
        "original_dimensions": f"{original_width}x{original_height}",
        "new_dimensions": f"{new_size[0]}x{new_size[1]}",
        "maintained_aspect_ratio": maintain_aspect,
    }


def _convert_image_format(
    file_path: str,
    output_format: str,
    output_path: Optional[str] = None,
    quality: int = 85,
) -> dict:
    """Convert an image to a different format."""
    try:
        from PIL import Image
    except ImportError:
        raise ImportError(
            "pillow package required for image tools. "
            "Install with: pip install tool-master[files]"
        )

    supported_formats = ['png', 'jpg', 'jpeg', 'webp', 'gif', 'bmp', 'tiff']
    output_format = output_format.lower()

    if output_format not in supported_formats:
        raise ValueError(f"Unsupported format: {output_format}. Supported: {supported_formats}")

    path = _validate_file_path(file_path)

    with Image.open(path) as img:
        # Handle format-specific conversions
        if output_format in ('jpg', 'jpeg') and img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')

        # Determine output path
        if output_path is None:
            ext = '.jpg' if output_format == 'jpeg' else f'.{output_format}'
            output_path = str(path.with_suffix(ext))

        output = Path(output_path).resolve()

        # Save with appropriate options
        save_kwargs = {}
        if output_format in ('jpg', 'jpeg', 'webp'):
            save_kwargs['quality'] = quality
        if output_format == 'png':
            save_kwargs['optimize'] = True

        img.save(output, **save_kwargs)

    return {
        "input_path": str(path),
        "output_path": str(output),
        "input_format": path.suffix[1:].upper(),
        "output_format": output_format.upper(),
        "quality": quality if output_format in ('jpg', 'jpeg', 'webp') else None,
    }


# =============================================================================
# TOOL DEFINITIONS
# =============================================================================

# Excel Tools
read_excel = Tool(
    name="read_excel",
    description="Read an Excel (.xlsx) file and return its contents as structured data with headers and rows.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the Excel file (.xlsx)",
            required=True,
        ),
        ToolParameter(
            name="sheet_name",
            type=ParameterType.STRING,
            description="Name of the sheet to read. Default: first/active sheet.",
            required=False,
        ),
        ToolParameter(
            name="max_rows",
            type=ParameterType.INTEGER,
            description="Maximum number of data rows to return (default: 1000).",
            required=False,
            default=1000,
        ),
    ],
    category="files",
    tags=["excel", "xlsx", "spreadsheet", "read", "files"],
).set_handler(_read_excel)


write_excel = Tool(
    name="write_excel",
    description="Write data to an Excel (.xlsx) file. Creates a new file or overwrites existing.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path for the output Excel file (.xlsx)",
            required=True,
        ),
        ToolParameter(
            name="data",
            type=ParameterType.ARRAY,
            description="2D array of data rows to write (e.g., [[1, 'a'], [2, 'b']])",
            required=True,
        ),
        ToolParameter(
            name="sheet_name",
            type=ParameterType.STRING,
            description="Name for the worksheet (default: 'Sheet1')",
            required=False,
            default="Sheet1",
        ),
        ToolParameter(
            name="headers",
            type=ParameterType.ARRAY,
            description="Optional header row (e.g., ['ID', 'Name'])",
            required=False,
        ),
    ],
    category="files",
    tags=["excel", "xlsx", "spreadsheet", "write", "files"],
).set_handler(_write_excel)


list_excel_sheets = Tool(
    name="list_excel_sheets",
    description="List all sheet names in an Excel workbook.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the Excel file (.xlsx)",
            required=True,
        ),
    ],
    category="files",
    tags=["excel", "xlsx", "spreadsheet", "sheets", "files"],
).set_handler(_list_excel_sheets)


read_excel_sheet_info = Tool(
    name="read_excel_sheet_info",
    description="Get detailed information about an Excel sheet (dimensions, headers, cell count).",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the Excel file (.xlsx)",
            required=True,
        ),
        ToolParameter(
            name="sheet_name",
            type=ParameterType.STRING,
            description="Name of the sheet. Default: active sheet.",
            required=False,
        ),
    ],
    category="files",
    tags=["excel", "xlsx", "spreadsheet", "info", "files"],
).set_handler(_read_excel_sheet_info)


# CSV Tools
read_csv = Tool(
    name="read_csv",
    description="Read a CSV file and return its contents. Auto-detects delimiter.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the CSV file",
            required=True,
        ),
        ToolParameter(
            name="delimiter",
            type=ParameterType.STRING,
            description="Column delimiter (default: auto-detect from content)",
            required=False,
        ),
        ToolParameter(
            name="max_rows",
            type=ParameterType.INTEGER,
            description="Maximum data rows to return (default: 1000)",
            required=False,
            default=1000,
        ),
        ToolParameter(
            name="encoding",
            type=ParameterType.STRING,
            description="File encoding (default: utf-8)",
            required=False,
            default="utf-8",
        ),
    ],
    category="files",
    tags=["csv", "read", "files", "data"],
).set_handler(_read_csv)


write_csv = Tool(
    name="write_csv",
    description="Write data to a CSV file.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path for the output CSV file",
            required=True,
        ),
        ToolParameter(
            name="data",
            type=ParameterType.ARRAY,
            description="2D array of data rows to write",
            required=True,
        ),
        ToolParameter(
            name="headers",
            type=ParameterType.ARRAY,
            description="Optional header row",
            required=False,
        ),
        ToolParameter(
            name="delimiter",
            type=ParameterType.STRING,
            description="Column delimiter (default: comma)",
            required=False,
            default=",",
        ),
        ToolParameter(
            name="encoding",
            type=ParameterType.STRING,
            description="File encoding (default: utf-8)",
            required=False,
            default="utf-8",
        ),
    ],
    category="files",
    tags=["csv", "write", "files", "data"],
).set_handler(_write_csv)


csv_to_excel = Tool(
    name="csv_to_excel",
    description="Convert a CSV file to Excel (.xlsx) format.",
    parameters=[
        ToolParameter(
            name="csv_path",
            type=ParameterType.STRING,
            description="Path to the source CSV file",
            required=True,
        ),
        ToolParameter(
            name="excel_path",
            type=ParameterType.STRING,
            description="Path for output Excel file (default: same name with .xlsx)",
            required=False,
        ),
        ToolParameter(
            name="sheet_name",
            type=ParameterType.STRING,
            description="Name for the worksheet (default: 'Sheet1')",
            required=False,
            default="Sheet1",
        ),
        ToolParameter(
            name="delimiter",
            type=ParameterType.STRING,
            description="CSV delimiter (default: auto-detect)",
            required=False,
        ),
    ],
    category="files",
    tags=["csv", "excel", "convert", "files"],
).set_handler(_csv_to_excel)


# JSON Tools
read_json = Tool(
    name="read_json",
    description="Read a JSON file and return its parsed contents.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the JSON file",
            required=True,
        ),
        ToolParameter(
            name="encoding",
            type=ParameterType.STRING,
            description="File encoding (default: utf-8)",
            required=False,
            default="utf-8",
        ),
    ],
    category="files",
    tags=["json", "read", "files", "data"],
).set_handler(_read_json)


write_json = Tool(
    name="write_json",
    description="Write data to a JSON file.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path for the output JSON file",
            required=True,
        ),
        ToolParameter(
            name="data",
            type=ParameterType.OBJECT,
            description="Data to write (dict, list, or primitive)",
            required=True,
        ),
        ToolParameter(
            name="pretty",
            type=ParameterType.BOOLEAN,
            description="Pretty-print with indentation (default: true)",
            required=False,
            default=True,
        ),
        ToolParameter(
            name="encoding",
            type=ParameterType.STRING,
            description="File encoding (default: utf-8)",
            required=False,
            default="utf-8",
        ),
    ],
    category="files",
    tags=["json", "write", "files", "data"],
).set_handler(_write_json)


validate_json = Tool(
    name="validate_json",
    description="Validate a JSON file and return its structure summary.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the JSON file to validate",
            required=True,
        ),
        ToolParameter(
            name="encoding",
            type=ParameterType.STRING,
            description="File encoding (default: utf-8)",
            required=False,
            default="utf-8",
        ),
    ],
    category="files",
    tags=["json", "validate", "files"],
).set_handler(_validate_json)


# PDF Tools
read_pdf_text = Tool(
    name="read_pdf_text",
    description="Extract text content from a PDF file, page by page.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the PDF file",
            required=True,
        ),
        ToolParameter(
            name="max_pages",
            type=ParameterType.INTEGER,
            description="Maximum pages to extract (default: all pages)",
            required=False,
        ),
        ToolParameter(
            name="page_numbers",
            type=ParameterType.ARRAY,
            description="Specific page numbers to extract (1-indexed)",
            required=False,
        ),
    ],
    category="files",
    tags=["pdf", "read", "text", "extract", "files"],
).set_handler(_read_pdf_text)


read_pdf_metadata = Tool(
    name="read_pdf_metadata",
    description="Get metadata from a PDF file (title, author, creation date, etc.).",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the PDF file",
            required=True,
        ),
    ],
    category="files",
    tags=["pdf", "metadata", "files"],
).set_handler(_read_pdf_metadata)


count_pdf_pages = Tool(
    name="count_pdf_pages",
    description="Get the page count of a PDF file.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the PDF file",
            required=True,
        ),
    ],
    category="files",
    tags=["pdf", "pages", "count", "files"],
).set_handler(_count_pdf_pages)


# PowerPoint Tools
read_pptx_text = Tool(
    name="read_pptx_text",
    description="Extract all text content from a PowerPoint file.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the PowerPoint file (.pptx)",
            required=True,
        ),
    ],
    category="files",
    tags=["powerpoint", "pptx", "text", "extract", "files"],
).set_handler(_read_pptx_text)


read_pptx_structure = Tool(
    name="read_pptx_structure",
    description="Get the structure of a PowerPoint file (slide titles, notes, shape counts).",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the PowerPoint file (.pptx)",
            required=True,
        ),
    ],
    category="files",
    tags=["powerpoint", "pptx", "structure", "files"],
).set_handler(_read_pptx_structure)


# Image Tools
read_image_metadata = Tool(
    name="read_image_metadata",
    description="Get metadata from an image file (dimensions, format, EXIF data).",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the image file",
            required=True,
        ),
    ],
    category="files",
    tags=["image", "metadata", "exif", "files"],
).set_handler(_read_image_metadata)


resize_image = Tool(
    name="resize_image",
    description="Resize an image to specified dimensions.",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the source image",
            required=True,
        ),
        ToolParameter(
            name="width",
            type=ParameterType.INTEGER,
            description="Target width in pixels",
            required=False,
        ),
        ToolParameter(
            name="height",
            type=ParameterType.INTEGER,
            description="Target height in pixels",
            required=False,
        ),
        ToolParameter(
            name="output_path",
            type=ParameterType.STRING,
            description="Output file path (default: input_resized.ext)",
            required=False,
        ),
        ToolParameter(
            name="maintain_aspect",
            type=ParameterType.BOOLEAN,
            description="Maintain aspect ratio (default: true)",
            required=False,
            default=True,
        ),
    ],
    category="files",
    tags=["image", "resize", "files"],
).set_handler(_resize_image)


convert_image_format = Tool(
    name="convert_image_format",
    description="Convert an image to a different format (png, jpg, webp, gif, bmp, tiff).",
    parameters=[
        ToolParameter(
            name="file_path",
            type=ParameterType.STRING,
            description="Path to the source image",
            required=True,
        ),
        ToolParameter(
            name="output_format",
            type=ParameterType.STRING,
            description="Target format: png, jpg, jpeg, webp, gif, bmp, tiff",
            required=True,
            enum=["png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff"],
        ),
        ToolParameter(
            name="output_path",
            type=ParameterType.STRING,
            description="Output file path (default: same name with new extension)",
            required=False,
        ),
        ToolParameter(
            name="quality",
            type=ParameterType.INTEGER,
            description="Quality for lossy formats like jpg/webp (1-100, default: 85)",
            required=False,
            default=85,
        ),
    ],
    category="files",
    tags=["image", "convert", "format", "files"],
).set_handler(_convert_image_format)
